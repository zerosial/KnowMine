"""
파서 서비스: Excel / PPT / PDF → Markdown 변환

핵심 원칙:
- 표(Table) 구조와 계층(Heading)을 완벽히 보존
- 시트명, 슬라이드 번호 등 메타 정보 포함
- AI가 컨텍스트를 파악하기 쉬운 구조적 Markdown 출력
"""
import io
import re
from pathlib import Path
from typing import Optional

import pandas as pd
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt


def detect_file_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    if ext in [".xlsx", ".xls", ".xlsm"]:
        return "excel"
    elif ext in [".pptx", ".ppt"]:
        return "ppt"
    elif ext == ".pdf":
        return "pdf"
    return "unknown"


# ─────────────────────────────────────────────
# Excel 파서
# ─────────────────────────────────────────────
def parse_excel(file_bytes: bytes, filename: str) -> str:
    """
    Excel 파일을 Markdown으로 변환.
    각 시트를 별도 섹션으로 분리, 표 구조 완전 보존.
    """
    md_sections = [f"# 📊 파일: {filename}\n"]

    try:
        xl = pd.ExcelFile(io.BytesIO(file_bytes))
        for sheet_name in xl.sheet_names:
            try:
                df = pd.read_excel(
                    io.BytesIO(file_bytes),
                    sheet_name=sheet_name,
                    header=0,
                    dtype=str,
                )
                # 완전히 빈 시트 skip
                if df.empty or df.dropna(how="all").empty:
                    continue

                df = df.fillna("")
                md_sections.append(f"\n## 📋 시트: {sheet_name}\n")

                # 병합 셀 처리: NaN → 이전 값 forward-fill
                df.ffill(inplace=True)

                # Markdown 테이블 생성
                md_table = _df_to_markdown(df)
                md_sections.append(md_table)

                # 간단한 통계 (숫자 컬럼만)
                numeric_df = df.apply(pd.to_numeric, errors="coerce")
                numeric_cols = numeric_df.columns[numeric_df.notna().any()].tolist()
                if numeric_cols:
                    stats_lines = ["\n> **요약 통계**"]
                    for col in numeric_cols[:5]:  # 최대 5개
                        col_data = numeric_df[col].dropna()
                        if not col_data.empty:
                            stats_lines.append(
                                f"> - `{col}`: 합계={col_data.sum():.2f}, 평균={col_data.mean():.2f}, 최대={col_data.max():.2f}"
                            )
                    md_sections.append("\n".join(stats_lines) + "\n")

            except Exception as e:
                md_sections.append(f"\n> ⚠️ 시트 `{sheet_name}` 파싱 오류: {str(e)}\n")

    except Exception as e:
        return f"# ❌ Excel 파싱 실패\n\n오류: {str(e)}"

    return "\n".join(md_sections)


def _df_to_markdown(df: pd.DataFrame) -> str:
    """DataFrame → Markdown 테이블 변환"""
    if df.empty:
        return "_데이터 없음_\n"

    # 컬럼 헤더
    headers = [str(col) for col in df.columns]
    header_row = "| " + " | ".join(headers) + " |"
    separator = "| " + " | ".join(["---"] * len(headers)) + " |"

    rows = []
    for _, row in df.iterrows():
        cells = [str(v).replace("|", "\\|").replace("\n", " ") for v in row]
        rows.append("| " + " | ".join(cells) + " |")

    return "\n".join([header_row, separator] + rows) + "\n"


# ─────────────────────────────────────────────
# PPT 파서
# ─────────────────────────────────────────────
def parse_pptx(file_bytes: bytes, filename: str) -> str:
    """
    PowerPoint 파일을 Markdown으로 변환.
    슬라이드별 제목/본문/표/노트 보존.
    """
    md_sections = [f"# 📑 파일: {filename}\n"]

    try:
        prs = Presentation(io.BytesIO(file_bytes))
        total_slides = len(prs.slides)
        md_sections.append(f"> 총 슬라이드: {total_slides}장\n")

        for idx, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"\n## 슬라이드 {idx}"]

            title = _get_slide_title(slide)
            if title:
                slide_parts.append(f"\n### {title}")

            for shape in slide.shapes:
                # 텍스트 프레임
                if shape.has_text_frame:
                    text = _extract_text_frame(shape.text_frame)
                    if text and text.strip() != title:
                        slide_parts.append(text)

                # 테이블
                if shape.has_table:
                    table_md = _pptx_table_to_markdown(shape.table)
                    slide_parts.append(table_md)

            # 발표자 노트
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n> 📝 노트: {notes_text}")

            md_sections.append("\n".join(slide_parts))

    except Exception as e:
        return f"# ❌ PPT 파싱 실패\n\n오류: {str(e)}"

    return "\n".join(md_sections)


def _get_slide_title(slide) -> Optional[str]:
    try:
        title_shape = slide.shapes.title
        if title_shape and title_shape.text:
            return title_shape.text.strip()
    except Exception:
        pass
    return None


def _extract_text_frame(text_frame) -> str:
    lines = []
    for para in text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level
        prefix = "  " * level + "-" if level > 0 else ""
        lines.append(f"{prefix} {text}" if prefix else text)
    return "\n".join(lines)


def _pptx_table_to_markdown(table) -> str:
    rows_data = []
    for row in table.rows:
        cells = [cell.text.replace("|", "\\|").replace("\n", " ").strip() for cell in row.cells]
        rows_data.append(cells)

    if not rows_data:
        return ""

    header = "| " + " | ".join(rows_data[0]) + " |"
    separator = "| " + " | ".join(["---"] * len(rows_data[0])) + " |"
    body_rows = ["| " + " | ".join(r) + " |" for r in rows_data[1:]]

    return "\n".join(["\n", header, separator] + body_rows) + "\n"


# ─────────────────────────────────────────────
# PDF 파서
# ─────────────────────────────────────────────
def parse_pdf(file_bytes: bytes, filename: str) -> str:
    """
    PDF 파일을 Markdown으로 변환.
    문서의 기본 폰트 크기(Base Size)를 동적으로 계산하여,
    폰트가 크거나 굵은(Bold) 텍스트를 범용적으로 제목(Heading)으로 인식합니다.
    """
    md_sections = [f"# 📄 파일: {filename}\n"]

    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        total_pages = len(doc)
        md_sections.append(f"> 총 페이지: {total_pages}페이지\n")

        # 1. 문서 전체를 스캔하여 기본 폰트 크기(Base Size) 도출
        font_counts = {}
        for page in doc:
            for block in page.get_text("dict").get("blocks", []):
                if block.get("type") == 0:
                    for line in block.get("lines", []):
                        for span in line.get("spans", []):
                            sz = round(span.get("size", 12))
                            txt = span.get("text", "").strip()
                            if txt:
                                font_counts[sz] = font_counts.get(sz, 0) + len(txt)
        
        # 가장 많이 쓰인 글자 크기를 본문(Body) 크기로 간주
        base_size = max(font_counts.items(), key=lambda x: x[1])[0] if font_counts else 12

        for page_num, page in enumerate(doc, start=1):
            page_md = [f"\n## 페이지 {page_num}"]

            # 텍스트 블록 추출
            blocks = page.get_text("dict").get("blocks", [])
            for block in blocks:
                if block.get("type") != 0:  # 0 = text block
                    continue
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue

                    text = "".join(s.get("text", "") for s in spans).strip()
                    if not text:
                        continue

                    # 범용 제목 인식 로직 (글자 크기, 볼드체 여부, 텍스트 길이 활용)
                    first_span = spans[0]
                    font_size = first_span.get("size", base_size)
                    font_name = first_span.get("font", "").lower()
                    flags = first_span.get("flags", 0)
                    
                    is_bold = ("bold" in font_name) or (flags & 16 != 0)
                    is_heading = False
                    
                    # 텍스트가 너무 길지 않으면서(제목의 특성)
                    if len(text) < 100:
                        if font_size >= base_size + 1.5:  # 본문보다 확연히 큼
                            is_heading = True
                        elif is_bold and font_size > base_size: # 본문보다 살짝 크고 볼드체
                            is_heading = True
                        elif is_bold and len(text) < 60: # 폰트 크기는 같지만 짧은 볼드체 강조 문구
                            is_heading = True
                    
                    if is_heading:
                        page_md.append(f"\n### {text}\n")
                    else:
                        page_md.append(text)

            # 테이블 추출 (PyMuPDF 내장)
            tables = page.find_tables()
            if tables.tables:
                for table in tables.tables:
                    try:
                        df = table.to_pandas()
                        if not df.empty:
                            page_md.append("\n**[표]**")
                            page_md.append(_df_to_markdown(df.astype(str).fillna("")))
                    except Exception:
                        pass

            md_sections.append("\n".join(page_md))
        doc.close()

    except Exception as e:
        return f"# ❌ PDF 파싱 실패\n\n오류: {str(e)}"

    return "\n".join(md_sections)


# ─────────────────────────────────────────────
# 통합 파서
# ─────────────────────────────────────────────
def parse_file(file_bytes: bytes, filename: str) -> tuple[str, str]:
    """
    파일 타입 자동 감지 후 파싱.
    Returns: (markdown_content, file_type)
    """
    file_type = detect_file_type(filename)

    if file_type == "excel":
        return parse_excel(file_bytes, filename), file_type
    elif file_type == "ppt":
        return parse_pptx(file_bytes, filename), file_type
    elif file_type == "pdf":
        return parse_pdf(file_bytes, filename), file_type
    else:
        # 알 수 없는 형식은 텍스트로 시도
        try:
            text = file_bytes.decode("utf-8", errors="ignore")
            return f"# 📁 파일: {filename}\n\n```\n{text}\n```\n", "unknown"
        except Exception:
            return f"# ❌ 지원하지 않는 파일 형식: {filename}", "unknown"
